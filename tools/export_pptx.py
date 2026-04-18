"""Export SECTIONS as a .pptx — companion to tools/import_pptx.py.

Each chapter gets a dark "lesson" cover slide; each case gets a two-column
slide (text left, image right if present). Speaker notes are written to the
notes placeholder so the exported deck survives re-import cleanly.

Usage:
    python3 tools/export_pptx.py --out /tmp/deck.pptx
    python3 tools/export_pptx.py --chapter 0 --out /tmp/ch0.pptx
"""
from __future__ import annotations
import argparse
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent))
from lint_deck import extract_sections  # noqa: E402

REPO_ROOT = Path(__file__).resolve().parent.parent
DEFAULT_HTML = REPO_ROOT / "index.html"

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dgm.color import RGBColor  # type: ignore  # noqa
except Exception:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
    except ImportError:
        print("ERROR: python-pptx not installed. pip3 install --user --break-system-packages python-pptx", file=sys.stderr)
        sys.exit(2)

# 16:9 default canvas is 13.33 × 7.5 inches.
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

ACCENT_RGB = {
    "teal":   RGBColor(0x0F, 0xB3, 0xA1),
    "purple": RGBColor(0x7C, 0x5C, 0xFF),
    "amber":  RGBColor(0xF5, 0xA5, 0x24),
    "rose":   RGBColor(0xF4, 0x3F, 0x7D),
}
BG = RGBColor(0x0B, 0x0B, 0x10)
FG = RGBColor(0xE8, 0xE8, 0xF0)
DIM = RGBColor(0xB8, 0xB8, 0xC8)


def _fill_bg(slide, rgb=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _add_textbox(slide, left, top, width, height, text, *,
                 size=18, bold=False, color=FG, line_height=1.2):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    lines = (text or "").split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_height
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def _is_local_image(img: str) -> Path | None:
    if not img or img.startswith(("IFRAME:", "MEDIA_CYCLER", "data:")):
        return None
    p = (REPO_ROOT / img) if not img.startswith("/") else Path(img)
    try:
        if p.exists() and p.suffix.lower() in (".png", ".jpg", ".jpeg", ".gif", ".webp"):
            return p
    except OSError:
        return None
    return None


def add_chapter_cover(prs: "Presentation", ch: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _fill_bg(slide)
    accent = ACCENT_RGB.get((ch.get("accent") or "").strip(), ACCENT_RGB["teal"])
    lesson = ch.get("lesson") or {}
    year = (ch.get("year") or "").strip()
    short = (lesson.get("short") or "").strip()
    header = year + ((" · " + short) if short and short != year else "")
    _add_textbox(slide, Inches(0.7), Inches(0.7), Inches(12), Inches(0.5),
                 header, size=14, bold=True, color=accent)
    title = (lesson.get("title") or "").strip()
    _add_textbox(slide, Inches(0.7), Inches(1.4), Inches(12), Inches(3),
                 title, size=44, bold=True, color=FG)
    tagline = (lesson.get("tagline") or "").strip()
    if tagline:
        _add_textbox(slide, Inches(0.7), Inches(5.2), Inches(12), Inches(1.6),
                     tagline, size=18, color=DIM)
    tags = (lesson.get("tags") or "").strip()
    if tags:
        _add_textbox(slide, Inches(0.7), Inches(6.7), Inches(12), Inches(0.5),
                     tags, size=12, color=accent)


def add_case_slide(prs: "Presentation", c: dict, accent: "RGBColor"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(slide)
    img_path = _is_local_image(c.get("img") or "")
    text_w = Inches(6.3) if img_path else Inches(12)

    _add_textbox(slide, Inches(0.7), Inches(0.7), text_w, Inches(1.2),
                 (c.get("title") or "").strip() or "(untitled)",
                 size=30, bold=True, color=FG)
    subtitle = (c.get("subtitle") or "").strip()
    if subtitle:
        _add_textbox(slide, Inches(0.7), Inches(1.9), text_w, Inches(0.9),
                     subtitle, size=16, color=DIM)

    bullets = [b for b in (c.get("bullets") or []) if isinstance(b, str) and b.strip()]
    if bullets:
        tb = slide.shapes.add_textbox(Inches(0.7), Inches(3.0), text_w, Inches(4.0))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.line_spacing = 1.3
            run = p.add_run()
            run.text = "• " + b.strip()
            run.font.size = Pt(16)
            run.font.color.rgb = FG

    if img_path:
        try:
            slide.shapes.add_picture(str(img_path), Inches(7.3), Inches(0.7),
                                     width=Inches(5.3), height=Inches(6.1))
        except Exception as e:
            print(f"    [pptx] image skipped ({img_path.name}): {e}", file=sys.stderr)

    notes = (c.get("notes") or "").strip()
    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    # Accent bar on the left edge.
    from pptx.shapes.autoshape import Shape  # noqa: F401
    bar = slide.shapes.add_shape(1, Inches(0.3), Inches(0.7), Inches(0.08), Inches(6.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()


def build(sections: list[dict], out_path: Path) -> tuple[int, int]:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    n_cases = 0
    for ch in sections:
        add_chapter_cover(prs, ch)
        accent = ACCENT_RGB.get((ch.get("accent") or "").strip(), ACCENT_RGB["teal"])
        for c in ch.get("cases") or []:
            add_case_slide(prs, c, accent)
            n_cases += 1
    prs.save(str(out_path))
    return len(sections), n_cases


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--html", type=Path, default=DEFAULT_HTML)
    ap.add_argument("--out", type=Path, required=True)
    ap.add_argument("--chapter", type=int, default=None)
    args = ap.parse_args()

    sections = extract_sections(args.html)
    if args.chapter is not None:
        if not (0 <= args.chapter < len(sections)):
            print(f"ERROR: chapter index out of range (0..{len(sections)-1})", file=sys.stderr)
            return 2
        sections = [sections[args.chapter]]

    n_ch, n_cases = build(sections, args.out)
    print(f"[done] Wrote {args.out} ({n_ch} chapters, {n_cases} cases, {n_ch + n_cases} slides)", file=sys.stderr)
    return 0


if __name__ == "__main__":
    sys.exit(main())
